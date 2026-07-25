from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
BODY_FONT = "PingFang SC"
MONO_FONT = "Menlo"
NAVY = PptRGBColor(20, 43, 64)
RED = PptRGBColor(190, 64, 54)
TEAL = PptRGBColor(35, 125, 125)
LIGHT = PptRGBColor(247, 248, 244)
MUTED = PptRGBColor(92, 102, 112)


MARKDOWN_DOCS = [
    (
        ROOT / "SCI_RESEARCH_VRP_LLM_DECISION.md",
        ROOT / "SCI_RESEARCH_VRP_LLM_DECISION.docx",
        ROOT / "SCI_RESEARCH_VRP_LLM_DECISION.pptx",
    ),
    (
        ROOT / "SCI_RESEARCH_FIRE_BLOCKAGE.md",
        ROOT / "SCI_RESEARCH_FIRE_BLOCKAGE.docx",
        ROOT / "SCI_RESEARCH_FIRE_BLOCKAGE.pptx",
    ),
]


def set_docx_font(run, size=None, bold=False, color=None, font_name=BODY_FONT):
    run.font.name = font_name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    run.bold = bold
    if size is not None:
        run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = RGBColor(*color)


def set_style_font(style, font_name=BODY_FONT, size=11):
    style.font.name = font_name
    style._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    style.font.size = Pt(size)


def add_code_block(document, lines):
    paragraph = document.add_paragraph()
    paragraph.paragraph_format.left_indent = Inches(0.25)
    paragraph.paragraph_format.space_after = Pt(6)
    run = paragraph.add_run("\n".join(lines))
    set_docx_font(run, size=9, font_name=MONO_FONT, color=(60, 66, 72))


def markdown_to_docx(markdown_path, output_path):
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)

    styles = document.styles
    set_style_font(styles["Normal"], size=10.5)
    set_style_font(styles["Heading 1"], size=18)
    set_style_font(styles["Heading 2"], size=15)
    set_style_font(styles["Heading 3"], size=12.5)
    set_style_font(styles["List Bullet"], size=10.5)

    in_code_block = False
    code_lines = []

    for raw_line in markdown_path.read_text(encoding="utf-8").splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()

        if stripped.startswith("```"):
            if in_code_block:
                add_code_block(document, code_lines)
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
                code_lines = []
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        if not stripped:
            continue

        if stripped.startswith("# "):
            paragraph = document.add_heading(stripped[2:], level=0)
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                set_docx_font(run, size=20, bold=True, color=(20, 43, 64))
            continue

        if stripped.startswith("## "):
            paragraph = document.add_heading(stripped[3:], level=1)
            for run in paragraph.runs:
                set_docx_font(run, size=16, bold=True, color=(20, 43, 64))
            continue

        if stripped.startswith("### "):
            paragraph = document.add_heading(stripped[4:], level=2)
            for run in paragraph.runs:
                set_docx_font(run, size=12.5, bold=True, color=(190, 64, 54))
            continue

        if stripped.startswith("- "):
            paragraph = document.add_paragraph(stripped[2:], style="List Bullet")
            paragraph.paragraph_format.space_after = Pt(2)
            for run in paragraph.runs:
                set_docx_font(run, size=10.5)
            continue

        paragraph = document.add_paragraph(stripped)
        paragraph.paragraph_format.first_line_indent = Pt(21)
        paragraph.paragraph_format.line_spacing = 1.18
        paragraph.paragraph_format.space_after = Pt(5)
        for run in paragraph.runs:
            set_docx_font(run, size=10.5)

    if code_lines:
        add_code_block(document, code_lines)

    document.save(output_path)


def add_textbox(slide, left, top, width, height, text, size=20, bold=False, color=NAVY, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = BODY_FONT
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide, title, section_label):
    add_textbox(slide, PptInches(0.55), PptInches(0.32), PptInches(11.4), PptInches(0.46), title, size=23, bold=True)
    add_textbox(slide, PptInches(11.55), PptInches(0.36), PptInches(1.25), PptInches(0.32), section_label, size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.55), PptInches(0.86), PptInches(12.2), PptInches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()


def add_bullets(slide, left, top, width, height, bullets, font_size=17):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = BODY_FONT
        paragraph.font.size = PptPt(font_size)
        paragraph.font.color.rgb = NAVY
        paragraph.space_after = PptPt(8)
    return box


def add_title_slide(prs, title, subtitle, label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(0.22), PptInches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()
    add_textbox(slide, PptInches(0.78), PptInches(1.55), PptInches(11.8), PptInches(1.35), title, size=34, bold=True)
    add_textbox(slide, PptInches(0.83), PptInches(3.08), PptInches(10.5), PptInches(0.58), subtitle, size=18, color=MUTED)
    add_textbox(slide, PptInches(0.83), PptInches(6.48), PptInches(5.6), PptInches(0.28), label, size=11, color=MUTED)


def add_bullet_slide(prs, title, bullets, label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_header(slide, title, label)
    add_bullets(slide, PptInches(0.78), PptInches(1.35), PptInches(11.7), PptInches(5.6), bullets)


def add_problem_content_slide(prs, title, problem, content, label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_header(slide, title, label)

    left_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.72), PptInches(1.35), PptInches(5.72), PptInches(5.25))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = PptRGBColor(255, 255, 255)
    left_box.line.color.rgb = PptRGBColor(220, 225, 228)
    right_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(6.85), PptInches(1.35), PptInches(5.72), PptInches(5.25))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = PptRGBColor(255, 255, 255)
    right_box.line.color.rgb = PptRGBColor(220, 225, 228)

    add_textbox(slide, PptInches(1.0), PptInches(1.65), PptInches(5.1), PptInches(0.35), "科学子问题", size=15, bold=True, color=RED)
    add_bullets(slide, PptInches(1.0), PptInches(2.15), PptInches(5.05), PptInches(3.9), problem, font_size=15)
    add_textbox(slide, PptInches(7.12), PptInches(1.65), PptInches(5.1), PptInches(0.35), "对应研究内容", size=15, bold=True, color=TEAL)
    add_bullets(slide, PptInches(7.12), PptInches(2.15), PptInches(5.05), PptInches(3.9), content, font_size=15)


def add_route_slide(prs, title, steps, label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_header(slide, title, label)
    x = 0.62
    y = 2.1
    width = 1.55
    gap = 0.25
    for index, step in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(width), PptInches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = PptRGBColor(255, 255, 255)
        box.line.color.rgb = PptRGBColor(205, 212, 216)
        text_frame = box.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = step
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = BODY_FONT
        paragraph.font.size = PptPt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = NAVY
        if index < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, PptInches(x + width + 0.02), PptInches(y + 0.35), PptInches(gap), PptInches(0.32))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RED
            arrow.line.fill.background()
        x += width + gap


def build_vrp_ppt(output_path):
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    add_title_slide(
        prs,
        "基于大模型决策的校园消防巡检路径规划",
        "LLM-assisted Risk-aware UAV Routing for Campus Fire Safety Inspection",
        "SCI 研究方向一",
    )
    add_bullet_slide(
        prs,
        "研究定位",
        [
            "传统 VRP 只描述几何访问顺序，难以表达消防巡检中的风险语义。",
            "大模型不直接替代优化器，而是生成风险权重、优先级、策略约束和解释。",
            "确定性优化器负责路径可行性，大模型负责上下文理解和动态决策。",
        ],
        "Background",
    )
    add_bullet_slide(
        prs,
        "核心科学问题",
        [
            "如何把校园消防语义、历史异常、时间上下文转化为可计算的巡检风险？",
            "如何让大模型参与路径规划决策，同时保证电量、禁飞区、任务时限等硬约束？",
            "如何利用异常反馈实现规划-巡检-识别-重规划闭环？",
        ],
        "Scientific Question",
    )
    add_problem_content_slide(
        prs,
        "子问题一：风险语义建模",
        ["消防巡检点不是普通坐标点。", "建筑功能、历史堵塞、时间段和消防设施都会影响巡检优先级。", "风险权重需要随场景动态变化。"],
        ["构建校园消防语义地图。", "建立节点多维属性：历史、人流、建筑、消防空间属性。", "设计风险权重计算模型。"],
        "Problem 1",
    )
    add_problem_content_slide(
        prs,
        "子问题二：大模型约束化决策",
        ["大模型直接输出路径容易产生不可行解。", "自然语言任务需求需要转换成 VRP 参数。", "开放式推理需要被约束校验。"],
        ["设计结构化提示和输出模板。", "让大模型输出权重、优先级和策略，不直接输出最终路径。", "建立约束校验器与优化器接口。"],
        "Problem 2",
    )
    add_problem_content_slide(
        prs,
        "子问题三：动态重规划",
        ["巡检过程中会出现疑似堵塞、电量下降、通信异常等动态事件。", "静态路径难以响应突发消防风险。", "需要平衡风险收益与路径代价。"],
        ["建立异常事件模型。", "设计异常触发的复检和重规划规则。", "研究局部重规划与高风险节点优先保留机制。"],
        "Problem 3",
    )
    add_route_slide(
        prs,
        "总体技术路线",
        ["消防语义地图", "风险知识库", "大模型决策", "约束校验", "VRP 求解", "无人机巡检", "异常反馈"],
        "Framework",
    )
    add_bullet_slide(
        prs,
        "可发表创新点",
        [
            "提出面向校园消防巡检的风险语义节点建模方法。",
            "构建大模型辅助的风险感知 VRP 决策框架。",
            "设计异常反馈驱动的动态重规划机制，实现闭环消防巡检。",
        ],
        "Contributions",
    )
    add_bullet_slide(
        prs,
        "实验验证设计",
        [
            "对比传统最短路径 VRP、固定权重风险 VRP 和无重规划策略。",
            "指标包括总距离、能耗、高风险节点覆盖率、异常响应时间和路径可行率。",
            "在不同电量、不同巡检时间窗口和不同异常密度下测试鲁棒性。",
        ],
        "Experiments",
    )
    prs.save(output_path)


def build_fire_ppt(output_path):
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    add_title_slide(
        prs,
        "面向无人机巡检的校园消防通道堵塞识别",
        "Relation-aware Transformer for UAV-based Campus Fire Lane Blockage Detection",
        "SCI 研究方向二",
    )
    add_bullet_slide(
        prs,
        "研究定位",
        [
            "消防堵塞识别不能等同于普通目标检测。",
            "检测到车辆或杂物，不一定意味着消防通道被堵塞。",
            "关键在于理解消防通道区域、障碍物和空间侵占关系。",
        ],
        "Background",
    )
    add_bullet_slide(
        prs,
        "核心科学问题",
        [
            "如何在无人机俯视/斜视图像中理解消防通道这种功能区域？",
            "如何鲁棒识别尺度小、类别开放、背景复杂的堵塞相关障碍物？",
            "如何从检测结果推理到消防通道是否被堵塞，以及堵塞等级？",
        ],
        "Scientific Question",
    )
    add_problem_content_slide(
        prs,
        "子问题一：消防通道区域理解",
        ["消防通道是功能区域，不只是道路纹理。", "普通道路、停车区和消防通道需要区分。", "无人机高度、视角和光照会造成识别不稳定。"],
        ["标注消防通道、消防栓禁占区、入口和登高操作面。", "构建消防通道电子围栏。", "利用地图先验增强区域识别。"],
        "Problem 1",
    )
    add_problem_content_slide(
        prs,
        "子问题二：堵塞相关障碍物识别",
        ["航拍目标尺度小、类别多、姿态变化大。", "临时堆放物和施工材料具有开放类别特征。", "正常停放物与通道占用物容易混淆。"],
        ["构建消防堵塞相关目标标注体系。", "设计多尺度 Transformer 特征表达。", "增强小目标、密集目标和开放类别识别能力。"],
        "Problem 2",
    )
    add_problem_content_slide(
        prs,
        "子问题三：堵塞关系推理",
        ["障碍物是否构成堵塞取决于位置、尺寸、类别和连通性影响。", "单纯检测框无法解释剩余通行能力。", "需要从对象识别转向空间关系理解。"],
        ["计算障碍物对通道的侵占比例。", "估计消防车剩余通行宽度。", "融合类别、持续时间和连通性生成堵塞等级。"],
        "Problem 3",
    )
    add_route_slide(
        prs,
        "FireLaneFormer 技术路线",
        ["无人机图像", "Transformer", "通道查询", "障碍物查询", "交互注意力", "侵占估计", "堵塞分级"],
        "Framework",
    )
    add_bullet_slide(
        prs,
        "可发表创新点",
        [
            "提出消防通道堵塞场景理解任务，而不是普通障碍物检测任务。",
            "设计通道-障碍物交互注意力，建模功能区域与目标之间的空间关系。",
            "建立可解释堵塞风险评分，输出侵占区域、剩余宽度和堵塞等级。",
        ],
        "Contributions",
    )
    add_bullet_slide(
        prs,
        "实验验证设计",
        [
            "对比 YOLO、DETR、Deformable DETR、Mask2Former、SegFormer 等方法。",
            "指标包括障碍物 mAP、通道 mIoU、堵塞等级 F1、侵占比例误差。",
            "通过消融实验验证通道分支、障碍物分支、交互注意力和风险评分模块。",
        ],
        "Experiments",
    )
    prs.save(output_path)


def main():
    for markdown_path, docx_path, _ in MARKDOWN_DOCS:
        markdown_to_docx(markdown_path, docx_path)
        print(f"wrote {docx_path.relative_to(ROOT)}")

    build_vrp_ppt(ROOT / "SCI_RESEARCH_VRP_LLM_DECISION.pptx")
    print("wrote SCI_RESEARCH_VRP_LLM_DECISION.pptx")
    build_fire_ppt(ROOT / "SCI_RESEARCH_FIRE_BLOCKAGE.pptx")
    print("wrote SCI_RESEARCH_FIRE_BLOCKAGE.pptx")


if __name__ == "__main__":
    main()